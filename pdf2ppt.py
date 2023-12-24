# -*- coding: utf-8 -*-
"""
Created on Thu Sep 14 16:11:31 2023

@author: Marloes Venema
@author: Christian Weinert
"""

import argparse
import io
import os

import pdf2image
import pptx

# Parse arguments for input and output files
parser = argparse.ArgumentParser(description="""
        Converts a PDF presentation into a PPT presentation.""",
                                 formatter_class=argparse.RawTextHelpFormatter)

parser.add_argument('--input', '-i', help='Specify the PDF input file', required=True)
parser.add_argument('--output', '-o', help='Specify the PPT output file')
parser.add_argument('--dpi', '-d', default=300, type=int, help='Specify the quality of the output file in DPI')

args = parser.parse_args()

# Get PDF input file name from args and check if file exists
pdf_name = args.input
if not os.path.exists(pdf_name):
    print("The specified PDF input file does not exist.")
    exit()

# Get PPT output file name from args
ppt_name = args.output

# If no PPT output file is specified, define output file based on input file
ppt_ext = ".pptx"
if ppt_name is None:
    ppt_name = os.path.splitext(pdf_name)[0] + ppt_ext

# Make sure the extension of the provided output file name is correct
else:
    ppt_name = os.path.splitext(ppt_name)[0] + ppt_ext

# Convert PDF with to high-res PNGs using multi-threading
images = pdf2image.convert_from_path(pdf_name, dpi=args.dpi, fmt='png', thread_count=os.cpu_count())

# Determine aspect ratio of PDF presentation
first_img = images[0]
ratio = first_img.width / float(first_img.height)

# Set up PPT presentation
prs = pptx.Presentation()

width = 10
prs.slide_width = pptx.util.Inches(width)
prs.slide_height = pptx.util.Inches(width / ratio)

# Add new slide for each PNG image
for img in images:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Convert image object to virtual file object
    tmpFile = io.BytesIO()
    img.save(tmpFile, "PNG")

    slide.shapes.add_picture(tmpFile, 0, 0, height=prs.slide_height)

prs.save(ppt_name)
