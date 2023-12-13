# -*- coding: utf-8 -*-
"""
Created on Thu Sep 14 16:11:31 2023

@author: Marloes Venema
"""

# import modules
from pdf2image import convert_from_path
import collections 
import collections.abc
from pptx import Presentation
import os
from pptx.util import Inches

# Change the pdf and ppt names accordingly
pdf_name = "test.pdf"
ppt_name = "test.pptx"

# Store Pdf with convert_from_path function
images = convert_from_path(pdf_name, 1000)

l = len(images)

for i in range(l):

 	# Save pages as images in the pdf
    images[i].save('slide'+ str(i+1) +'.png', 'PNG')
    
# Now convert to presentation
prs = Presentation()
# The width and height below correspond to 16:9 ratio
# For 4:3 ratio, you may want to use 10 and 7.5
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
for i in range(l):
    # print("converting slide" + str(i+1) + " now")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_path = 'slide'+ str(i+1) +'.png'
    left = Inches(0)
    top = Inches(0)
    height = Inches(5.625)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    
prs.save(ppt_name)
os.startfile(ppt_name)
